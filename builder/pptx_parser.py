"""Parse a .pptx Support Notes content template into structured slot data.

The PPTX template consists of:
  Slide 1: Instructions — ignored.
  Slide 2: Main content table (typically 10 rows × 7 cols).  Column 6 holds
           human-readable section labels that describe the structure.
  Slides 3+: One picture shape per slide — these are speaker headshot images.

Content table cell mapping:
  [0, 0]  section1.quote
  [1, 0]  section1.editorial_html
  [1, 4]  section1 speaker  (name + title below "[Headshot here]" placeholder)
  [2, 0]  section1.left_articles  (one hyperlinked paragraph per article)
  [3, 0]  section2.quote
  [4, 0]  section2 speaker  (name + title below "[Headshot here]" placeholder)
  [4, 1]  section2.editorial_html
  [5, 0]  section2.left_articles  (one hyperlinked paragraph per article)
  [6, 0]  webinar contact block → webinar_series_url extracted from hyperlinks
  [6, 3]  latest_release.components  (one hyperlinked paragraph per component)
  row 8+  footer — skipped

All hyperlinks are resolved through the slide's relationship dictionary.
Picture shapes on slides other than the content slide are returned as
EmlAttachment objects so the caller can offer them as speaker headshot
candidates (the exact same type used by the .eml parser).
"""

from __future__ import annotations

import io
import re
from pathlib import Path
from typing import IO, List, Optional, Tuple, Union

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from builder.models import (
    Article,
    EditorialSection,
    EmlAttachment,
    ReleaseComponent,
    ReleaseHighlight,
    Speaker,
    SupportNotesData,
)

# "[Headshot here]" / "[Headshot here ]" placeholder in speaker cells
_HEADSHOT_PLACEHOLDER_RE = re.compile(r'^\[.*?headshot.*?\]', re.IGNORECASE)

# Generic support-centre root URL — not treated as the webinar series URL
_SUPPORT_ROOT = "https://support.sw.siemens.com/en-US"


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def parse(
    source: Union[str, Path, bytes, IO[bytes]],
) -> Tuple[SupportNotesData, List[EmlAttachment]]:
    """Parse a .pptx file and return ``(data, headshot_images)``.

    ``data`` has all extractable content fields populated directly from the
    PPTX table.  ``product``, ``language``, ``year``, and ``month`` are left
    at their defaults (empty / 0) for the user to fill in via the Streamlit UI,
    since the PPTX has no equivalent of the .eml Subject header.

    ``headshot_images`` is a list of ``EmlAttachment`` objects containing the
    embedded picture shapes from slides other than the content slide.  The
    caller can display these as thumbnails and let the user pick which one to
    composite and upload for each speaker section.
    """
    prs = _load_pptx(source)
    content_slide_idx, content_slide, table_shape = _find_content_slide(prs)

    rel_map = {
        r_id: rel.target_ref
        for r_id, rel in content_slide.part.rels.items()
        if hasattr(rel, "target_ref")
    }

    table = table_shape.table
    data = _extract_data(table, rel_map)
    images = _collect_images(prs, content_slide_idx)
    return data, images


# ---------------------------------------------------------------------------
# Data extraction
# ---------------------------------------------------------------------------

def _extract_data(table, rel_map: dict) -> SupportNotesData:
    """Map the PPTX table rows/columns to SupportNotesData fields."""
    # Section 1: editorial left, headshot right
    s1_editorial = _cell_text(table, 1, 0)
    s1_speaker   = _parse_speaker(_cell_text(table, 1, 4))
    s1_articles  = _cell_articles(_safe_cell(table, 2, 0), rel_map)

    # Section 2: headshot left, editorial right
    s2_speaker   = _parse_speaker(_cell_text(table, 4, 0))
    s2_editorial = _cell_text(table, 4, 1)
    s2_articles  = _cell_articles(_safe_cell(table, 5, 0), rel_map)

    # Section 3: webinar contact block + latest release
    webinar_cell  = _safe_cell(table, 6, 0)
    release_cell  = _safe_cell(table, 6, 3)
    webinar_url   = _cell_webinar_series_url(webinar_cell, rel_map) if webinar_cell else ""
    release_comps = _cell_release_components(release_cell, rel_map) if release_cell else []

    return SupportNotesData(
        section1=EditorialSection(
            quote=_cell_text(table, 0, 0),
            editorial_html=s1_editorial,
            speaker=s1_speaker,
            left_articles=s1_articles,
        ),
        section2=EditorialSection(
            quote=_cell_text(table, 3, 0),
            editorial_html=s2_editorial,
            speaker=s2_speaker,
            left_articles=s2_articles,
        ),
        webinar_series_url=webinar_url,
        latest_release=ReleaseHighlight(components=release_comps),
    )


# ---------------------------------------------------------------------------
# Cell helpers
# ---------------------------------------------------------------------------

def _safe_cell(table, row: int, col: int):
    """Return the cell at (row, col) or None if out of bounds."""
    try:
        return table.rows[row].cells[col]
    except IndexError:
        return None


def _cell_text(table, row: int, col: int) -> str:
    """Return the plain text of a cell, stripped of whitespace."""
    cell = _safe_cell(table, row, col)
    if cell is None or cell.text_frame is None:
        return ""
    return cell.text_frame.text.strip()


def _parse_speaker(cell_text: str) -> Speaker:
    """Extract Speaker name and title from a speaker-block cell.

    The cell contains a "[Headshot here]" placeholder followed by blank lines,
    then the name on one line and the title on the next, e.g.::

        [Headshot here ]


        Sandhya Nageswaren
        Applications Engineer
    """
    lines = [l.strip() for l in cell_text.splitlines()]
    # Drop the headshot placeholder line and any blank lines
    content = [l for l in lines if l and not _HEADSHOT_PLACEHOLDER_RE.match(l)]
    name  = content[0] if content else ""
    title = content[1] if len(content) > 1 else "Applications Engineer"
    return Speaker(name=name, title=title)


def _cell_articles(cell, rel_map: dict) -> List[Article]:
    """Return one Article per hyperlinked paragraph in the cell."""
    if cell is None:
        return []
    return [
        Article(label=text, url=url)
        for text, url in _hyperlinked_paragraphs(cell, rel_map)
    ]


def _cell_release_components(cell, rel_map: dict) -> List[ReleaseComponent]:
    """Return one ReleaseComponent per hyperlinked paragraph in the cell."""
    if cell is None:
        return []
    return [
        ReleaseComponent(label=text, url=url)
        for text, url in _hyperlinked_paragraphs(cell, rel_map)
    ]


def _cell_webinar_series_url(cell, rel_map: dict) -> str:
    """Return the webinar series URL from the contact/webinar cell.

    The cell contains several hyperlinks; the generic support-centre root is
    excluded and the first remaining URL is returned.
    """
    for _text, url in _hyperlinked_paragraphs(cell, rel_map):
        if url.rstrip("/") != _SUPPORT_ROOT:
            return url
    return ""


# ---------------------------------------------------------------------------
# Hyperlink extraction
# ---------------------------------------------------------------------------

def _hyperlinked_paragraphs(cell, rel_map: dict) -> List[Tuple[str, str]]:
    """Return ``[(paragraph_text, url), ...]`` for paragraphs with a hyperlink.

    Each paragraph is checked for any run that carries an ``a:hlinkClick``
    element.  If found, the entire paragraph's text is paired with the
    resolved URL.  Paragraphs whose runs all share one r:id (e.g. a component
    name split across multiple formatting runs) produce a single entry with the
    combined paragraph text.
    """
    results = []
    if cell is None or cell.text_frame is None:
        return results
    for para in cell.text_frame.paragraphs:
        para_text = para.text.strip()
        if not para_text:
            continue
        url = _first_hyperlink_url(para, rel_map)
        if url:
            results.append((para_text, url))
    return results


def _first_hyperlink_url(para, rel_map: dict) -> Optional[str]:
    """Return the first hyperlink URL found in the paragraph's runs, or None."""
    for run in para.runs:
        rpr = run._r.find(qn("a:rPr"))
        if rpr is None:
            continue
        hlink = rpr.find(qn("a:hlinkClick"))
        if hlink is None:
            continue
        r_id = hlink.get(qn("r:id"))
        if r_id and r_id in rel_map:
            return rel_map[r_id]
    return None


# ---------------------------------------------------------------------------
# Image extraction
# ---------------------------------------------------------------------------

def _collect_images(prs: Presentation, content_slide_idx: int) -> List[EmlAttachment]:
    """Collect picture shapes from all slides except the content slide.

    Returns one EmlAttachment per picture shape.  Small images (likely logos
    or decorative icons — under 50 KB) are skipped.
    """
    images: List[EmlAttachment] = []
    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx == content_slide_idx:
            continue
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                blob = shape.image.blob
                if len(blob) < 50_000:
                    continue
                ext = shape.image.ext
                ct = shape.image.content_type
                images.append(EmlAttachment(
                    filename=f"headshot_slide{slide_idx + 1}.{ext}",
                    content_type=ct,
                    content_id=None,
                    bytes=blob,
                ))
            except Exception:
                pass
    return images


# ---------------------------------------------------------------------------
# PPTX loading
# ---------------------------------------------------------------------------

def _load_pptx(source: Union[str, Path, bytes, IO[bytes]]) -> Presentation:
    if isinstance(source, (str, Path)):
        return Presentation(str(source))
    if isinstance(source, bytes):
        return Presentation(io.BytesIO(source))
    if hasattr(source, "read"):
        return Presentation(source)
    raise TypeError(f"Unsupported source type for .pptx parse: {type(source)!r}")


def _find_content_slide(prs: Presentation):
    """Return ``(slide_index, slide, table_shape)`` for the slide with the largest table."""
    best: Optional[tuple] = None  # (cell_count, slide_idx, slide, shape)
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            size = len(shape.table.rows) * len(shape.table.columns)
            if best is None or size > best[0]:
                best = (size, slide_idx, slide, shape)
    if best is None:
        raise ValueError(
            "No table found in the PPTX — cannot parse as a Support Notes template."
        )
    _, slide_idx, slide, shape = best
    return slide_idx, slide, shape
