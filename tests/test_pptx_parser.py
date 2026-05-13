"""Tests for builder.pptx_parser.

Exercises the heuristic extraction from a PPTX content table using a
synthetic in-memory presentation rather than the real sample file.
"""

from __future__ import annotations

import io

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from builder import pptx_parser
from builder.models import EmlAttachment


# ---------------------------------------------------------------------------
# Helpers — build a minimal PPTX that matches the expected table structure
# ---------------------------------------------------------------------------

def _make_test_pptx() -> bytes:
    """Build a two-slide PPTX whose second slide has the standard content table."""
    prs = Presentation()

    # Slide 1: instructions (blank layout, no table)
    blank_layout = prs.slide_layouts[6]
    prs.slides.add_slide(blank_layout)

    # Slide 2: content table (10 rows × 7 cols)
    slide = prs.slides.add_slide(blank_layout)
    rows, cols = 10, 7
    left = top = Inches(0.5)
    width = Inches(9)
    height = Inches(6)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table

    def _set(r, c, text):
        tbl.rows[r].cells[c].text_frame.text = text

    _set(0, 0, "Section 1 headline goes here")
    _set(0, 6, "Section 1: Quote / Headline from Support team")

    _set(1, 0, "Editorial body for section 1.")
    _set(1, 4, "[Headshot here]\n\n\n\nAlice Example\nApplications Engineer")
    _set(1, 6, "Section 2: editorial (left), headshot (right)")

    # Row 2 — resource links (no hyperlinks in this synthetic build, so no articles)
    _set(2, 6, "Section 3: Resources")

    _set(3, 0, "Section 2 headline")
    _set(3, 6, "Section 4: Quote / Headline from Support team")

    _set(4, 0, "[Headshot here]\n\n\n\nBob Example\nSenior Engineer")
    _set(4, 1, "Editorial body for section 2.")
    _set(4, 6, "Section 5: Headshot (left), editorial (right)")

    _set(5, 6, "Section 6: Resources")

    _set(6, 0, "For support, contact us via the Support Center.")
    _set(6, 6, "Section 7")

    _set(8, 0, "Footer text — ignored")
    _set(8, 6, "Default Siemens footer")

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------

@pytest.fixture(scope="module")
def pptx_bytes():
    return _make_test_pptx()


def test_section1_quote(pptx_bytes):
    data, _ = pptx_parser.parse(pptx_bytes)
    assert data.section1.quote == "Section 1 headline goes here"


def test_section1_editorial(pptx_bytes):
    data, _ = pptx_parser.parse(pptx_bytes)
    assert data.section1.editorial_html == "Editorial body for section 1."


def test_section1_speaker_name(pptx_bytes):
    data, _ = pptx_parser.parse(pptx_bytes)
    assert data.section1.speaker.name == "Alice Example"


def test_section1_speaker_title(pptx_bytes):
    data, _ = pptx_parser.parse(pptx_bytes)
    assert data.section1.speaker.title == "Applications Engineer"


def test_section2_quote(pptx_bytes):
    data, _ = pptx_parser.parse(pptx_bytes)
    assert data.section2.quote == "Section 2 headline"


def test_section2_editorial(pptx_bytes):
    data, _ = pptx_parser.parse(pptx_bytes)
    assert data.section2.editorial_html == "Editorial body for section 2."


def test_section2_speaker(pptx_bytes):
    data, _ = pptx_parser.parse(pptx_bytes)
    assert data.section2.speaker.name == "Bob Example"
    assert data.section2.speaker.title == "Senior Engineer"


def test_no_articles_when_no_hyperlinks(pptx_bytes):
    data, _ = pptx_parser.parse(pptx_bytes)
    assert data.section1.left_articles == []
    assert data.section2.left_articles == []


def test_no_release_components_when_no_hyperlinks(pptx_bytes):
    data, _ = pptx_parser.parse(pptx_bytes)
    assert data.latest_release.components == []


def test_no_images_in_synthetic_pptx(pptx_bytes):
    _, images = pptx_parser.parse(pptx_bytes)
    assert images == []


def test_parse_from_path(tmp_path):
    """parse() accepts a file path as well as bytes."""
    p = tmp_path / "test.pptx"
    p.write_bytes(_make_test_pptx())
    data, images = pptx_parser.parse(p)
    assert data.section1.quote == "Section 1 headline goes here"


def test_parse_from_filelike():
    """parse() accepts a binary file-like object (e.g. Streamlit UploadedFile)."""
    buf = io.BytesIO(_make_test_pptx())
    data, _ = pptx_parser.parse(buf)
    assert data.section1.speaker.name == "Alice Example"


def test_headshot_placeholder_stripped(pptx_bytes):
    """The [Headshot here] placeholder text must NOT appear in speaker name."""
    data, _ = pptx_parser.parse(pptx_bytes)
    assert "Headshot" not in data.section1.speaker.name
    assert "Headshot" not in data.section2.speaker.name


def test_footer_not_extracted(pptx_bytes):
    """Row 8 (footer) must be ignored — the content should not appear anywhere."""
    data, _ = pptx_parser.parse(pptx_bytes)
    d = data.to_dict()
    import json
    serialised = json.dumps(d)
    assert "Footer text" not in serialised
